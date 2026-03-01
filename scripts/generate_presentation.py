#!/usr/bin/env python3
"""
Generate EngageIQ Flight Deck PowerPoint presentation.
Run: pip install python-pptx && python scripts/generate_presentation.py
Output: engageiq-flight/slides/EngageIQ_Flight_Deck.pptx
Add screenshots to slides/screenshots/ (warning.png, critical.png, etc.) to auto-insert.
"""
from __future__ import annotations

import os
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("Install: pip install python-pptx")
    raise

OUTPUT_DIR = Path(__file__).parent.parent / "slides"
SCREENSHOTS_DIR = OUTPUT_DIR / "screenshots"


def _add_slide(prs: Presentation, layout_idx: int = 1) -> "Slide":
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def _set_title(slide, text: str) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = text


def _add_body(slide, text: str) -> None:
    for sp in slide.placeholders:
        if sp.placeholder_format.idx == 1 and hasattr(sp, "text_frame"):
            sp.text_frame.text = text
            return
    # Fallback: add textbox
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)


def _add_image_if_exists(slide, path_or_stem: Path | str, left: float = 4, top: float = 2, width: float = 4.5) -> bool:
    """Add image if exists. path_or_stem can be full path or stem (tries .png, .jpg)."""
    p = Path(path_or_stem)
    if p.suffix:
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(left), Inches(top), width=Inches(width))
            return True
        return False
    for ext in (".png", ".jpg", ".jpeg", ".jpe", ".webp"):
        candidate = SCREENSHOTS_DIR / f"{p.name}{ext}"
        if candidate.exists():
            slide.shapes.add_picture(str(candidate), Inches(left), Inches(top), width=Inches(width))
            return True
    return False


def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    SCREENSHOTS_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ---- 1. Title / Introduction ----
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "EngageIQ Flight Deck"
    s.placeholders[1].text = "Real-Time Pilot Behavioral Safety Monitoring\nGoogle DeepMind × InstaLILY On-Device AI Hackathon"

    # ---- 2. Real-Life Problem ----
    s = _add_slide(prs)
    _set_title(s, "Real-Life Problem")
    _add_body(
        s,
        "• Pilot fatigue causes ~15–20% of aviation accidents\n"
        "• Biometric video (face) cannot leave the aircraft — privacy & security regulations\n"
        "• Cloud round-trip latency at cruise (satellite): 600–800ms — too slow for &lt;500ms alerts\n"
        "• Oceanic routes: zero connectivity for hours — system must work offline\n"
        "• Continuous camera feed economics: per-call API pricing is impractical"
    )

    # ---- 3. Solving Use Case ----
    s = _add_slide(prs)
    _set_title(s, "Solving Use Case")
    _add_body(
        s,
        "• Two camera feeds + cockpit data → live Pilot Safety Index (PSI, 0–100)\n"
        "• Escalating safety interventions: NOMINAL → MONITOR → CAUTION → WARNING → CRITICAL\n"
        "• Fully on-device inference via Ollama (Gemma 3n, FunctionGemma) — no cloud\n"
        "• Autonomous ActionAgent decides: trigger_alert, notify_copilot, suggest_rest_protocol, etc.\n"
        "• Demo: upload drowsy video or inject scenarios; watch PSI drop and alerts escalate"
    )

    # ---- 4. Components ----
    s = _add_slide(prs)
    _set_title(s, "Components")
    _add_body(
        s,
        "• PerceptionAgent (Gemma 3n) — face frame → pilot_state (perclos, yawn, head_drop, gaze)\n"
        "• ActionAgent (FunctionGemma 270M) — pilot_state → function calls\n"
        "• LandingAgent (Gemma 3n) — external/belly frame → landing_report (bounces, score)\n"
        "• PSI Engine — composite score from fatigue, attention, procedural, landing\n"
        "• Cockpit Error Classifier — 5 error types (wrong_button, omission, commission, reversal, sequence_error)\n"
        "• React HUD — PSI ring, signals, agent feed, landing report, alert banner"
    )

    # ---- 5. Architecture ----
    s = _add_slide(prs)
    _set_title(s, "Architecture")
    _add_body(
        s,
        "Camera 1 (Face) → OpenCV → PerceptionAgent (Gemma 3n) → pilot_state JSON\n"
        "Camera 2 (External) → OpenCV → LandingAgent (Gemma 3n) → landing_report JSON\n"
        "X-Plane / Replay → sim_bridge → phase + cockpit_errors\n"
        "pilot_state + cockpit_errors → ActionAgent (FunctionGemma) → function calls\n"
        "All outputs → PSI Engine → WebSocket (2Hz) → React HUD"
    )

    # ---- 6–10. What Each Part Does ----
    parts = [
        ("PerceptionAgent", "Gemma 3n E4B\n\n• Input: face frame (Camera 1) + flight phase\n• Output: pilot_state JSON (state, perclos_est, yawn, head_drop, gaze_off_instruments)\n• Runs every 2 seconds\n• LoRA fine-tuned on NTHU drowsy driver dataset"),
        ("ActionAgent", "FunctionGemma 270M (fine-tuned)\n\n• Input: pilot_state stream + session context\n• Output: function call strings (trigger_alert, notify_copilot, suggest_rest_protocol, log_fatigue_event, request_atc_advisory)\n• Autonomous safety loop — decides when to escalate"),
        ("LandingAgent", "Gemma 3n E4B\n\n• Input: external/belly frame (Camera 2) during approach/landing\n• Output: landing_report (bounce_count, on_centerline, contact_type, score)\n• Scores: greaser 100, firm 80, hard 60, 4+ bounces 10"),
        ("PSI Engine", "Composite 0–100 score\n\n• Fatigue: perclos, yawn count, CRITICAL events (−40/−15/−30 pts)\n• Attention: head drop, gaze off (−30 pts)\n• Procedural: cockpit errors (−30 pts)\n• Landing: penalty if last score &lt;40 (−10 pts)\n• Alert bands: 85+ NOMINAL, 70–84 MONITOR, 55–69 CAUTION, 35–54 WARNING, 0–34 CRITICAL"),
        ("React HUD", "Live dashboard\n\n• PSI ring (animated score)\n• Signal grid (yawns, head drops, gaze %)\n• Agent feed (FunctionGemma decisions)\n• Landing report (bounces, score)\n• Alert banner (escalating colors)\n• PSI timeline (Recharts)\n• Video panel (upload Cam 1/2), Simulator panel (inject scenarios)"),
    ]
    for title, body in parts:
        s = _add_slide(prs)
        _set_title(s, f"What It Does: {title}")
        _add_body(s, body)

    # ---- 11. Warning Scenarios ----
    s = _add_slide(prs)
    _set_title(s, "Warning Scenarios")
    _add_body(
        s,
        "• Drowsy — elevated perclos, gaze drift → trigger_alert(WARNING), suggest_rest_protocol(15)\n"
        "• Fatigue onset — perclos rising → log_fatigue_event, trigger_alert(WARNING)\n"
        "• Gaze drift — gaze off instruments → trigger_alert(WARNING)\n"
        "• Yawn series — 4+ yawns/10m → suggest_rest_protocol(20)"
    )
    _add_image_if_exists(s, "drowsy", left=0.5, top=3.5, width=4)
    _add_image_if_exists(s, "warning_ui", left=5, top=3.5, width=4.5)

    # ---- 12. Critical Scenarios ----
    s = _add_slide(prs)
    _set_title(s, "Critical Scenarios")
    _add_body(
        s,
        "• Critical — micro-sleep risk, head drop → notify_copilot, trigger_alert(CRITICAL)\n"
        "• Micro-sleep — eyes closed &gt;3s → notify_copilot, request_atc_advisory\n"
        "• Head drop severe — eyelids drooping → notify_copilot\n"
        "• Cockpit reversal — throttle opposite → trigger_alert(CRITICAL)"
    )
    _add_image_if_exists(s, "critical", left=0.5, top=3.5, width=4)
    _add_image_if_exists(s, "critical_ui", left=5, top=3.5, width=4.5)

    # ---- 13. Takeoff / Landing Scenarios ----
    s = _add_slide(prs)
    _set_title(s, "Takeoff & Landing Scenarios")
    _add_body(
        s,
        "• Takeoff smooth — nominal climb\n"
        "• Approach nominal — stable approach\n"
        "• Landing greaser — score 100, 1 contact\n"
        "• Landing hard — score 60, G&gt;1.8\n"
        "• Landing bounce — score 10, 4+ bounces, runway excursion risk\n"
        "• Go-around — initiated"
    )
    _add_image_if_exists(s, "landing_bounce", left=0.5, top=3.5, width=4)
    _add_image_if_exists(s, "landing_ui", left=5, top=3.5, width=4.5)

    # ---- 14. Cockpit Error Scenarios ----
    s = _add_slide(prs)
    _set_title(s, "Cockpit Error Scenarios")
    _add_body(
        s,
        "• Gear omission — gear not down by 1000 ft AGL (−8 pts)\n"
        "• Wrong button — wrong flap setting (−20 pts)\n"
        "• Reversal — throttle opposite (−20 pts)\n"
        "• Sequence error — checklist out of order (−2 pts)\n"
        "• Cockpit multi — fatigue + procedural errors combined"
    )
    _add_image_if_exists(s, "cockpit_errors", left=0.5, top=3.5, width=4)
    _add_image_if_exists(s, "cockpit_ui", left=5, top=3.5, width=4.5)

    # ---- 15. Key Claims ----
    s = _add_slide(prs)
    _set_title(s, "Key Claims")
    _add_body(
        s,
        "• All inference on-device; no cloud API calls\n"
        "• Pilot biometric video never leaves the aircraft\n"
        "• Alert latency &lt;500ms; cloud impossible at cruise\n"
        "• Fully offline capable for oceanic routes\n"
        "• Gemma 3n reasons over frames; not hardcoded thresholds\n"
        "• FunctionGemma autonomously decides alerts\n"
        "• Fine-tuned models improve accuracy\n"
        "• Prototype; real deployment requires DO-178C, FAA/EASA"
    )

    # ---- 16. Demo Flow ----
    s = _add_slide(prs)
    _set_title(s, "Demo Flow")
    _add_body(
        s,
        "1. Cold HUD — PSI 95, NOMINAL\n"
        "2. Inject drowsy / upload NTHU video → PSI drops, WARNING\n"
        "3. Inject critical → CRITICAL, notify_copilot\n"
        "4. Inject landing_bounce → score 10, runway excursion risk\n"
        "5. Inject cockpit_errors → procedural deductions\n"
        "6. Turn off WiFi — system keeps running"
    )

    # ---- 17. Thank You ----
    s = _add_slide(prs)
    _set_title(s, "Thank You")
    _add_body(s, "EngageIQ Flight Deck\nReal-time pilot safety — fully on-device.\n\nQuestions?")

    out = OUTPUT_DIR / "EngageIQ_Flight_Deck.pptx"
    try:
        prs.save(str(out))
    except PermissionError:
        out = OUTPUT_DIR / "EngageIQ_Flight_Deck_with_images.pptx"
        prs.save(str(out))
        print("(Original file may be open; saved as _with_images.pptx)")
    print(f"Saved: {out}")
    print(f"Add screenshots to {SCREENSHOTS_DIR} (warning.png, critical.png, landing.png, cockpit.png, etc.) and re-run to insert them.")


if __name__ == "__main__":
    main()
